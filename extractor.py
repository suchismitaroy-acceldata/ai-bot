import os
from pypdf import PdfReader
from docx import Document
from pptx import Presentation

MAX_FILE_SIZE_MB = 15

def extract_text(file_path):
    try:
        size_mb = os.path.getsize(file_path) / (1024 * 1024)

        # 🔥 Skip large files
        if size_mb > MAX_FILE_SIZE_MB:
            print("⏭ Skipping large file:", file_path)
            return ""

        # 🔥 PDF
        if file_path.endswith(".pdf"):
            reader = PdfReader(file_path)

            text = []
            for i, page in enumerate(reader.pages):
                if i > 10:  # 🔥 LIMIT PAGES
                    break
                text.append(page.extract_text() or "")

            return " ".join(text)

        # 🔥 DOCX
        elif file_path.endswith(".docx"):
            doc = Document(file_path)
            return "\n".join([p.text for p in doc.paragraphs[:100]])

        # 🔥 PPTX
        elif file_path.endswith(".pptx"):
            prs = Presentation(file_path)
            text = []
            for slide in prs.slides[:20]:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)
            return "\n".join(text)

        # 🔥 TXT
        elif file_path.endswith(".txt"):
            with open(file_path, "r", encoding="utf-8") as f:
                return f.read()

        return ""

    except Exception as e:
        print("❌ Error reading:", file_path, e)
        return ""